from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def criar_apresentacao_facss():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cores Oficiais FACSS
    BG_COLOR = RGBColor(11, 19, 17)        # #0b1311
    CARD_BG = RGBColor(16, 28, 24)         # #101c18
    GREEN_PRIMARY = RGBColor(43, 168, 112) # #2ba870
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94a3b8

    def aplicar_fundo_escuro(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # -------------------------------------------------------------------------
    # SLIDE 1: CAPA
    # -------------------------------------------------------------------------
    blank_layout = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_escuro(s1)

    tx_box = s1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "FACSS CRM & OPERAÇÕES"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = GREEN_PRIMARY

    p1 = tf.add_paragraph()
    p1.text = "Plataforma Integrada de Gestão B2B, Suporte Operacional e Telemetria"
    p1.font.size = Pt(22)
    p1.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "\nDesenvolvido por Flávio Alves | Python • Flask • Firebase • Cloud Deploy"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED

    # -------------------------------------------------------------------------
    # SLIDE 2: O DESAFIO VS A SOLUÇÃO
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_escuro(s2)

    title_box = s2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    tf2 = title_box.text_frame
    p = tf2.paragraphs[0]
    p.text = "O Desafio vs. A Solução"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = GREEN_PRIMARY

    itens = [
        ("Desafio Anterior", "Dados espalhados e controle manual em planilhas sem rastreabilidade."),
        ("Solução Centralizada", "Repositório NoSQL em nuvem (Firebase) para dados em tempo real."),
        ("Pipeline de Vendas", "Visibilidade do funil B2B através de quadro Kanban intuitivo."),
        ("Retenção Preditiva", "Alertas automáticos de inatividade de clientes e controle de Health Score.")
    ]

    left = Inches(1.0)
    top = Inches(2.2)
    width = Inches(5.3)
    height = Inches(2.0)

    for i, (titulo, desc) in enumerate(itens):
        r = i // 2
        c = i % 2
        col_left = left + c * Inches(5.8)
        row_top = top + r * Inches(2.3)

        shape = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_left, row_top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = GREEN_PRIMARY

        tf_card = shape.text_frame
        tf_card.word_wrap = True
        p_t = tf_card.paragraphs[0]
        p_t.text = titulo
        p_t.font.bold = True
        p_t.font.size = Pt(18)
        p_t.font.color.rgb = GREEN_PRIMARY

        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # SLIDE 3: FUNCIONALIDADES-CHAVE
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_escuro(s3)

    t_box3 = s3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    p3 = t_box3.text_frame.paragraphs[0]
    p3.text = "Funcionalidades Principais"
    p3.font.size = Pt(32)
    p3.font.bold = True
    p3.font.color.rgb = GREEN_PRIMARY

    funcs = [
        "📊 Painel Geral: Métricas executivas, saúde de carteira e gráficos com Chart.js.",
        "🎯 Pipeline Kanban: Gestão de oportunidades B2B em estágios customizáveis.",
        "🛠️ Ordens de Serviço: Chamados técnicos, criticidade e emissão de relatórios em PDF.",
        "📧 Notificações SMTP: Disparo de e-mails de confirmação e boas-vindas via SSL.",
        "📡 Módulos Operacionais: Controle de rastreadores, SIM Cards e vínculo com frotas."
    ]

    func_box = s3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.8))
    tf_f = func_box.text_frame
    tf_f.word_wrap = True

    for i, f_text in enumerate(funcs):
        p_f = tf_f.paragraphs[0] if i == 0 else tf_f.add_paragraph()
        p_f.text = f_text + "\n"
        p_f.font.size = Pt(18)
        p_f.font.color.rgb = TEXT_WHITE

    # -------------------------------------------------------------------------
    # SLIDE 4: ARQUITETURA TÉCNICA
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    aplicar_fundo_escuro(s4)

    t_box4 = s4.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.0))
    p4 = t_box4.text_frame.paragraphs[0]
    p4.text = "Arquitetura & Tecnologias"
    p4.font.size = Pt(32)
    p4.font.bold = True
    p4.font.color.rgb = GREEN_PRIMARY

    techs = [
        ("Backend", "Python 3.11, Flask Framework, Jinja2 Templates"),
        ("Banco de Dados", "Google Cloud Firebase Firestore (NoSQL)"),
        ("Frontend", "Bootstrap 5.3 (Dark Theme), Chart.js, HTML5/CSS3"),
        ("Infraestrutura", "Render Cloud Platform, Git/GitHub, SMTP SSL")
    ]

    for i, (camada, tech) in enumerate(techs):
        row_top = Inches(2.0) + i * Inches(1.2)
        
        shape = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), row_top, Inches(11.3), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = GREEN_PRIMARY

        tf_t = shape.text_frame
        p_c = tf_t.paragraphs[0]
        p_c.text = f"{camada}: "
        p_c.font.bold = True
        p_c.font.size = Pt(16)
        p_c.font.color.rgb = GREEN_PRIMARY
        
        run = p_c.add_run()
        run.text = tech
        run.font.color.rgb = TEXT_WHITE
        run.font.bold = False

    # Salvar arquivo
    prs.save("FACSS_CRM_Apresentacao.pptx")
    print("✅ Apresentação 'FACSS_CRM_Apresentacao.pptx' gerada com sucesso!")

if __name__ == "__main__":
    criar_apresentacao_facss()