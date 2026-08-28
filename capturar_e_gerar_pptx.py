import os
import time
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# -----------------------------------------------------------------------------
# PASSO 1: ROBÔ QUE LOGA NO CRM E TIRA OS PRINTS REAIS
# -----------------------------------------------------------------------------
def capturar_telas_crm():
    print("🚀 Iniciando navegador para capturar as telas do CRM...")
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(viewport={'width': 1600, 'height': 900})
        page = context.new_page()

        # 1. Login no CRM
        page.goto('https://portal-facss.onrender.com/login')
        page.fill('input[name="email"]', 'flavio.alves@facss.com.br')
        page.fill('input[name="senha"]', 'Facss2026')
        page.click('button[type="submit"]')
        page.wait_for_timeout(3000)

        # 2. Captura dos Módulos Principais
        telas = [
            ('/', '01_dashboard.png'),
            ('/clientes', '02_clientes.png'),
            ('/pipeline', '03_pipeline.png'),
            ('/os', '04_ordens_servico.png'),
            ('/modulos', '05_modulos.png')
        ]

        for rota, nome_arquivo in telas:
            url = f'https://portal-facss.onrender.com{rota}'
            page.goto(url)
            page.wait_for_timeout(2000)
            page.screenshot(path=nome_arquivo, full_page=False)
            print(f"📸 Screenshot salva: {nome_arquivo}")

        browser.close()
    print("✅ Todas as telas foram capturadas com sucesso!")

# -----------------------------------------------------------------------------
# PASSO 2: MONTAGEM DO POWERPOINT COM AS IMAGENS EMBUTIDAS
# -----------------------------------------------------------------------------
def gerar_powerpoint_com_prints():
    print("📊 Montando apresentação do PowerPoint com as telas...")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_COLOR = RGBColor(11, 19, 17)
    GREEN_PRIMARY = RGBColor(43, 168, 112)
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)

    def aplicar_fundo(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    blank_layout = prs.slide_layouts[6]

    # SLIDE 1: CAPA
    s1 = prs.slides.add_slide(blank_layout)
    aplicar_fundo(s1)
    
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(3.0))
    p0 = tb.text_frame.paragraphs[0]
    p0.text = "FACSS CRM & OPERAÇÕES"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = GREEN_PRIMARY

    p1 = tb.text_frame.add_paragraph()
    p1.text = "Demonstração Visual da Intranet Enterprise B2B"
    p1.font.size = Pt(22)
    p1.font.color.rgb = TEXT_WHITE

    p2 = tb.text_frame.add_paragraph()
    p2.text = "\nDesenvolvido por Flávio Alves | Python • Flask • Firebase"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED

    # SLIDES DE CADA TELA
    slides_data = [
        ("Painel Geral & Indicadores (KPIs)", "01_dashboard.png", "Métricas em tempo real, saúde da carteira e carga de chamados."),
        ("Gestão de Carteira B2B", "02_clientes.png", "Listagem de clientes ativos, responsáveis e status de retenção."),
        ("Pipeline Kanban de Vendas", "03_pipeline.png", "Acompanhamento visual de oportunidades por estágios de prospecção."),
        ("Central de Ordens de Serviço (OS)", "04_ordens_servico.png", "Controle de chamados com níveis de criticidade e emissão de laudos."),
        ("Módulos Operacionais & Telemetria", "05_modulos.png", "Inventário de rastreadores, chips SIM e vínculo com clientes.")
    ]

    for titulo, imagem, sub in slides_data:
        s = prs.slides.add_slide(blank_layout)
        aplicar_fundo(s)

        # Título do Slide
        tb_title = s.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
        pt = tb_title.text_frame.paragraphs[0]
        pt.text = titulo
        pt.font.size = Pt(24)
        pt.font.bold = True
        pt.font.color.rgb = GREEN_PRIMARY

        ps = tb_title.text_frame.add_paragraph()
        ps.text = sub
        ps.font.size = Pt(13)
        ps.font.color.rgb = TEXT_MUTED

        # Inserção da Screenshot
        if os.path.exists(imagem):
            s.shapes.add_picture(imagem, Inches(0.8), Inches(1.3), Inches(11.733), Inches(5.6))

    prs.save("FACSS_CRM_Apresentacao_Oficial.pptx")
    print("🎉 Apresentação 'FACSS_CRM_Apresentacao_Oficial.pptx' criada com as fotos do sistema!")

if __name__ == '__main__':
    try:
        capturar_telas_crm()
        gerar_powerpoint_com_prints()
    except Exception as e:
        print(f"❌ Erro durante a automação: {e}")